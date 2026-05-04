"""
Generate Release Analytics presentation.
Run: .venv/Scripts/python generate_ppt.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG        = RGBColor(0x0D, 0x0D, 0x1F)   # near-black navy
C_CARD      = RGBColor(0x16, 0x16, 0x2E)   # card background
C_PURPLE    = RGBColor(0x63, 0x66, 0xF1)   # indigo accent
C_PURPLE_LT = RGBColor(0x81, 0x8C, 0xF8)   # lighter purple
C_GREEN     = RGBColor(0x34, 0xD3, 0x99)   # emerald
C_RED       = RGBColor(0xF8, 0x71, 0x71)   # red (before)
C_BLUE      = RGBColor(0x60, 0xA5, 0xFA)   # blue
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY      = RGBColor(0x94, 0xA3, 0xB8)
C_DARK_CARD = RGBColor(0x1E, 0x1E, 0x3A)

W = Inches(13.33)   # widescreen 16:9
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank


# ── Low-level helpers ─────────────────────────────────────────────────────────

def solid_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def box(slide, x, y, w, h, color, alpha=None):
    from pptx.oxml.ns import qn
    shape = slide.shapes.add_shape(1, x, y, w, h)   # MSO_SHAPE_TYPE.RECTANGLE = 1
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(slide, text, x, y, w, h,
        size=18, bold=False, italic=False,
        color=C_WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(x, y, w, h)
    txb.word_wrap = wrap
    tf  = txb.text_frame
    tf.word_wrap = wrap
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def accent_bar(slide, x, y, h=Inches(0.04), color=C_PURPLE):
    """Thin horizontal accent line."""
    box(slide, x, y, W - x * 2, h, color)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═════════════════════════════════════════════════════════════════════════════
def slide_title(prs):
    sl = prs.slides.add_slide(BLANK)
    solid_bg(sl, C_BG)

    # Left purple stripe
    box(sl, Inches(0), Inches(0), Inches(0.18), H, C_PURPLE)

    # Big gradient-style header block
    box(sl, Inches(0.18), Inches(1.8), Inches(13.15), Inches(3.4), C_DARK_CARD)
    box(sl, Inches(0.18), Inches(1.8), Inches(0.1),   Inches(3.4), C_PURPLE_LT)

    txt(sl, "Release Analytics", Inches(0.6), Inches(2.0), Inches(12), Inches(1.1),
        size=52, bold=True, color=C_WHITE)
    txt(sl, "Internal Platform — Engineering & Delivery Review",
        Inches(0.6), Inches(3.1), Inches(12), Inches(0.65),
        size=22, color=C_PURPLE_LT)
    txt(sl, "How we went from scattered spreadsheets to a live, multi-team intelligence platform",
        Inches(0.6), Inches(3.75), Inches(12), Inches(0.55),
        size=14, italic=True, color=C_GREY)

    # Bottom meta
    txt(sl, "Expense on Demand  ·  April 2026  ·  Internal",
        Inches(0.6), Inches(6.8), Inches(12), Inches(0.4),
        size=11, color=C_GREY)

    # Decorative dots top-right
    for i, c in enumerate([C_PURPLE, C_GREEN, C_BLUE]):
        box(sl, Inches(11.5 + i * 0.5), Inches(0.3), Inches(0.28), Inches(0.28), c)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE JOURNEY  (v1 → v2 → v3)
# ═════════════════════════════════════════════════════════════════════════════
def slide_journey(prs):
    sl = prs.slides.add_slide(BLANK)
    solid_bg(sl, C_BG)
    box(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.85), C_DARK_CARD)

    txt(sl, "The Journey", Inches(0.5), Inches(0.12), Inches(8), Inches(0.65),
        size=28, bold=True, color=C_WHITE)
    txt(sl, "Three phases. One platform.", Inches(0.5), Inches(0.55), Inches(8), Inches(0.35),
        size=13, color=C_GREY)

    # Version cards
    versions = [
        {
            "label":  "v1",
            "title":  "Delivery Metrics",
            "period": "Phase 1",
            "color":  C_PURPLE,
            "points": [
                "Sprint tracking & burndown",
                "Team velocity trends",
                "Work item age & priority",
                "ADO data pipeline — live refresh",
            ],
        },
        {
            "label":  "v2",
            "title":  "Quality & Capacity",
            "period": "Phase 2",
            "color":  C_BLUE,
            "points": [
                "Bug MTTC, hotspot analysis",
                "QA health & test coverage",
                "Per-developer utilisation",
                "Estimation accuracy tracking",
            ],
        },
        {
            "label":  "v3",
            "title":  "Platform & Intelligence",
            "period": "Phase 3 — Now",
            "color":  C_GREEN,
            "points": [
                "In-house project management",
                "Epic / Release / Feature model",
                "Auth + RBAC for team rollout",
                "AI assistant layer",
            ],
        },
    ]

    card_w = Inches(3.9)
    gap    = Inches(0.3)
    start_x = Inches(0.45)

    for i, v in enumerate(versions):
        cx = start_x + i * (card_w + gap)
        cy = Inches(1.05)
        ch = Inches(5.85)

        box(sl, cx, cy, card_w, ch, C_DARK_CARD)
        box(sl, cx, cy, card_w, Inches(0.06), v["color"])

        # Version badge
        badge = box(sl, cx + Inches(0.2), cy + Inches(0.15), Inches(0.6), Inches(0.42), v["color"])

        txt(sl, v["label"], cx + Inches(0.2), cy + Inches(0.15), Inches(0.6), Inches(0.42),
            size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        txt(sl, v["title"],  cx + Inches(0.9), cy + Inches(0.15), Inches(2.8), Inches(0.42),
            size=17, bold=True, color=C_WHITE)
        txt(sl, v["period"], cx + Inches(0.9), cy + Inches(0.55), Inches(2.8), Inches(0.3),
            size=11, color=v["color"])

        # Divider
        box(sl, cx + Inches(0.2), cy + Inches(0.95), card_w - Inches(0.4), Inches(0.02),
            RGBColor(0x30, 0x30, 0x50))

        for j, pt in enumerate(v["points"]):
            py = cy + Inches(1.1) + j * Inches(0.75)
            box(sl, cx + Inches(0.22), py + Inches(0.18), Inches(0.1), Inches(0.1), v["color"])
            txt(sl, pt, cx + Inches(0.42), py, card_w - Inches(0.6), Inches(0.65),
                size=13, color=C_GREY)

    # Arrow connectors between cards
    for i in range(2):
        ax = start_x + (i + 1) * (card_w + gap) - Inches(0.25)
        txt(sl, "→", ax, Inches(3.5), Inches(0.28), Inches(0.4),
            size=22, bold=True, color=C_PURPLE_LT, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — LIVE DEMO
# ═════════════════════════════════════════════════════════════════════════════
def slide_demo(prs):
    sl = prs.slides.add_slide(BLANK)
    solid_bg(sl, C_BG)

    box(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.85), C_DARK_CARD)
    txt(sl, "Live Demo", Inches(0.5), Inches(0.12), Inches(8), Inches(0.65),
        size=28, bold=True, color=C_WHITE)

    # Big centre block
    box(sl, Inches(1.5), Inches(1.3), Inches(10.33), Inches(4.8), C_DARK_CARD)
    box(sl, Inches(1.5), Inches(1.3), Inches(10.33), Inches(0.06), C_GREEN)

    txt(sl, "🖥", Inches(5.8), Inches(2.2), Inches(1.8), Inches(1.5),
        size=60, align=PP_ALIGN.CENTER)

    txt(sl, "127.0.0.1:8050",
        Inches(1.5), Inches(3.5), Inches(10.33), Inches(0.7),
        size=28, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)

    txt(sl, "Release Analytics — Live Dashboard",
        Inches(1.5), Inches(4.2), Inches(10.33), Inches(0.5),
        size=16, color=C_GREY, align=PP_ALIGN.CENTER)

    # Nav items
    nav_items = ["Home", "Summary", "Planning", "Items", "Teams", "Assistant", "Platform"]
    nw = Inches(10.33) / len(nav_items)
    for i, nav in enumerate(nav_items):
        nx = Inches(1.5) + i * nw
        c  = C_PURPLE if nav == "Platform" else RGBColor(0x25, 0x25, 0x45)
        box(sl, nx + Inches(0.05), Inches(5.15), nw - Inches(0.1), Inches(0.55), c)
        txt(sl, nav, nx + Inches(0.05), Inches(5.18), nw - Inches(0.1), Inches(0.5),
            size=11, color=C_WHITE, align=PP_ALIGN.CENTER)

    txt(sl, "Authentication active  ·  Role-based access  ·  PostgreSQL backend  ·  Real ADO data",
        Inches(1.5), Inches(5.9), Inches(10.33), Inches(0.4),
        size=11, italic=True, color=C_GREY, align=PP_ALIGN.CENTER)


# ═════════════════════════════════════════════════════════════════════════════
# SLIDES 4–10 — ONE PER AREA
# ═════════════════════════════════════════════════════════════════════════════

AREAS = [
    {
        "icon":   "📊",
        "title":  "Release Status",
        "before": "Scattered across ADO, emails, and spreadsheets",
        "after":  "Single live dashboard — 8 pages, real time",
        "color":  C_PURPLE,
        "detail": [
            "All work items from ADO synced via direct DB pipeline",
            "Summary page: KPIs, bug matrix, backlog trend, burn rate by priority",
            "Release Outlook: burn rate, scope changes, delivery risk",
            "No manual assembly — refresh the page, data is current",
        ],
        "metric_label": "Pages of live analytics",
        "metric_value": "8",
    },
    {
        "icon":   "📈",
        "title":  "Reporting",
        "before": "Hours assembling weekly manual reports",
        "after":  "Always-on — refresh = done",
        "color":  C_BLUE,
        "detail": [
            "Data pipeline runs on demand — no scheduled jobs needed",
            "Every chart regenerates in seconds from live ADO data",
            "Week-on-week trends maintained automatically",
            "Zero manual data extraction or formatting",
        ],
        "metric_label": "Time to produce a report",
        "metric_value": "~3s",
    },
    {
        "icon":   "👥",
        "title":  "Capacity Insight",
        "before": "Impossible — no assigned vs built-by split",
        "after":  "Per-developer utilisation & accuracy tracked",
        "color":  C_GREEN,
        "detail": [
            "Separate main_developer field tracked independently of assigned_to",
            "Per-person utilisation: actual hours vs available hours",
            "Estimation accuracy: original estimate vs completed work",
            "Team and individual drill-down with iteration filtering",
        ],
        "metric_label": "Data accuracy improvement",
        "metric_value": "57%",
    },
    {
        "icon":   "🐛",
        "title":  "Bug Analytics",
        "before": "Ad-hoc ADO queries, no trend view",
        "after":  "MTTC, hotspots, source analysis, matrices",
        "color":  C_RED,
        "detail": [
            "Mean Time to Close (MTTC) tracked per team and priority",
            "Bug source analysis — where bugs originate in the pipeline",
            "Week-on-week movement matrix: opened vs closed vs net",
            "Priority hotspot chart — where critical bugs accumulate",
        ],
        "metric_label": "Bug metrics tracked",
        "metric_value": "6+",
    },
    {
        "icon":   "🎯",
        "title":  "Data Accuracy",
        "before": "57% team mismatch (assigned vs builder)",
        "after":  "main_dev_team fix — Dev work counted correctly",
        "color":  C_PURPLE_LT,
        "detail": [
            "Root cause: assigned_to often QA (ticket owner), not the builder",
            "Fix: separate main_dev_team column derived from main_developer field",
            "Capacity page uses main_dev_team — other pages use assigned_to correctly",
            "3,295 / 5,763 work items were previously attributed to the wrong team",
        ],
        "metric_label": "Items correctly re-attributed",
        "metric_value": "3,295",
    },
    {
        "icon":   "⏱",
        "title":  "Hour Counting",
        "before": "Double-counting parent + child estimates",
        "after":  "Self-correcting exclusion rule — zero false totals",
        "color":  RGBColor(0xFB, 0x92, 0x3C),
        "detail": [
            "Root/Middle Enhancement containers with completed_work=0 excluded",
            "These are new-style parent containers — their child Tasks hold the hours",
            "Historical data (parents with hours) unaffected — no data loss",
            "Rule fires automatically — no manual correction needed",
        ],
        "metric_label": "False hours removed",
        "metric_value": "Auto",
    },
    {
        "icon":   "🤖",
        "title":  "AI Integration",
        "before": "Zero",
        "after":  "AI Assistant page live in the dashboard",
        "color":  C_GREEN,
        "detail": [
            "Natural language interface to query ADO data",
            "Ask about metrics, trends, team performance, blockers",
            "Integrated into the same authenticated platform",
            "No context switching — analytics and AI in one place",
        ],
        "metric_label": "AI pages in platform",
        "metric_value": "1",
    },
]


def slide_area(prs, area):
    sl = prs.slides.add_slide(BLANK)
    solid_bg(sl, C_BG)

    c = area["color"]

    # Top bar
    box(sl, Inches(0), Inches(0), Inches(13.33), Inches(1.1), C_DARK_CARD)
    box(sl, Inches(0), Inches(0), Inches(13.33), Inches(0.05), c)

    # Icon + title
    txt(sl, area["icon"], Inches(0.3), Inches(0.1), Inches(0.85), Inches(0.85),
        size=34, align=PP_ALIGN.CENTER)
    txt(sl, area["title"], Inches(1.1), Inches(0.1), Inches(9), Inches(0.6),
        size=30, bold=True, color=C_WHITE)
    txt(sl, "Before → After", Inches(1.1), Inches(0.65), Inches(9), Inches(0.35),
        size=12, color=C_GREY)

    # Metric badge (top right)
    box(sl, Inches(10.8), Inches(0.15), Inches(2.3), Inches(0.8), c)
    txt(sl, area["metric_value"], Inches(10.8), Inches(0.1), Inches(2.3), Inches(0.5),
        size=26, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(sl, area["metric_label"], Inches(10.8), Inches(0.58), Inches(2.3), Inches(0.35),
        size=9, color=C_WHITE, align=PP_ALIGN.CENTER)

    # BEFORE card
    box(sl, Inches(0.4), Inches(1.25), Inches(5.9), Inches(1.35), RGBColor(0x2A, 0x10, 0x10))
    box(sl, Inches(0.4), Inches(1.25), Inches(5.9), Inches(0.05), C_RED)
    txt(sl, "BEFORE", Inches(0.55), Inches(1.3), Inches(2), Inches(0.35),
        size=10, bold=True, color=C_RED)
    txt(sl, area["before"], Inches(0.55), Inches(1.62), Inches(5.6), Inches(0.85),
        size=14, color=RGBColor(0xFF, 0xB0, 0xB0))

    # Arrow
    txt(sl, "→", Inches(6.45), Inches(1.55), Inches(0.5), Inches(0.7),
        size=28, bold=True, color=c, align=PP_ALIGN.CENTER)

    # AFTER card
    box(sl, Inches(7.05), Inches(1.25), Inches(5.9), Inches(1.35), RGBColor(0x0A, 0x20, 0x18))
    box(sl, Inches(7.05), Inches(1.25), Inches(5.9), Inches(0.05), C_GREEN)
    txt(sl, "AFTER", Inches(7.2), Inches(1.3), Inches(2), Inches(0.35),
        size=10, bold=True, color=C_GREEN)
    txt(sl, area["after"], Inches(7.2), Inches(1.62), Inches(5.6), Inches(0.85),
        size=14, color=RGBColor(0xA0, 0xFF, 0xD0))

    # Detail points
    txt(sl, "What changed", Inches(0.4), Inches(2.75), Inches(6), Inches(0.4),
        size=12, bold=True, color=c)
    box(sl, Inches(0.4), Inches(3.12), Inches(12.55), Inches(0.02),
        RGBColor(0x25, 0x25, 0x45))

    for i, pt in enumerate(area["detail"]):
        py = Inches(3.2) + i * Inches(0.82)
        box(sl, Inches(0.4), py + Inches(0.25), Inches(0.12), Inches(0.12), c)
        txt(sl, pt, Inches(0.65), py, Inches(12.2), Inches(0.75),
            size=14, color=C_GREY)

    # Bottom accent
    box(sl, Inches(0), Inches(7.2), Inches(13.33), Inches(0.3), C_DARK_CARD)
    txt(sl, "Expense on Demand  ·  Release Analytics",
        Inches(0.4), Inches(7.22), Inches(12), Inches(0.25),
        size=9, color=RGBColor(0x45, 0x45, 0x65))


# ═════════════════════════════════════════════════════════════════════════════
# BUILD
# ═════════════════════════════════════════════════════════════════════════════

slide_title(prs)
slide_journey(prs)
slide_demo(prs)
for area in AREAS:
    slide_area(prs, area)

out = r"c:\Python\Release\Release_Analytics_Presentation.pptx"
prs.save(out)
print(f"Saved: {out}")
print(f"Total slides: {len(prs.slides)}")
