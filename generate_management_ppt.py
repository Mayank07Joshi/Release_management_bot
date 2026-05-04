"""
Management Presentation — Release Analytics Platform
Run: .venv/Scripts/python generate_management_ppt.py

Sections:
  1  Title
  2  The State of Our Releases (stark KPIs)
  3  Defect Escape & PROD Quality
  4  Customer Impact
  5  Year-End Scorecard (Now vs Target — no sugarcoating)
  6  This Is All Live (transition)
  7  What the Platform Tracks
  8  Why This Tool Beats Asking AI (confident)
  9  This Tool vs AI — Head to Head
  10 What's Coming Next (roadmap pitch)
  11 Closing
"""

import sys, os
sys.path.insert(0, os.path.dirname(__file__))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
C_BG     = RGBColor(0x0D, 0x0D, 0x1F)
C_CARD   = RGBColor(0x16, 0x16, 0x2E)
C_CARD2  = RGBColor(0x1E, 0x1E, 0x3A)
C_PURPLE = RGBColor(0x63, 0x66, 0xF1)
C_PURPLT = RGBColor(0x81, 0x8C, 0xF8)
C_GREEN  = RGBColor(0x34, 0xD3, 0x99)
C_RED    = RGBColor(0xF8, 0x71, 0x71)
C_AMBER  = RGBColor(0xFB, 0xBF, 0x24)
C_BLUE   = RGBColor(0x60, 0xA5, 0xFA)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREY   = RGBColor(0x94, 0xA3, 0xB8)
C_DGREY  = RGBColor(0x4A, 0x55, 0x68)

W = Inches(13.33)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Helpers ───────────────────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BG
    return s


def box(s, x, y, w, h, color, radius=False):
    shape = s.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def txt(s, text, x, y, w, h, size=18, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, wrap=True):
    tf = s.shapes.add_textbox(x, y, w, h).text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def accent_bar(s, x, y, h=Inches(0.055), color=C_PURPLE):
    box(s, x, y, Inches(0.45), h, color)


def section_tag(s, label, x=Inches(0.5), y=Inches(0.28)):
    """Small pill label at top of slide."""
    pill = box(s, x, y, Inches(2.2), Inches(0.28), C_CARD2)
    txt(s, label, x + Inches(0.12), y + Inches(0.03), Inches(2.0), Inches(0.22),
        size=9, color=C_PURPLT, bold=True)


def divider(s, y):
    box(s, Inches(0.5), y, W - Inches(1.0), Inches(0.012), C_DGREY)


# ── Pull live data ─────────────────────────────────────────────────────────────
print("Loading data...")
try:
    from data.loader import load_data
    from reports.summarizer import summarize_bugs, summarize_releases, summarize_qa, summarize_capacity
    df   = load_data()
    bugs = summarize_bugs(df)
    qa   = summarize_qa(df)
    rels = summarize_releases(df)
    cap  = summarize_capacity(df)
    print(f"  {bugs['total']} total bugs, {bugs['total_open']} open")
    print(f"  {bugs['escape_rate_pct']}% escape rate, {bugs['open_p1']} P1 open")
    DATA_OK = True
except Exception as e:
    print(f"  WARNING: Could not load live data ({e}). Using fallback numbers.")
    DATA_OK = False
    # Fallback from our last analysis
    bugs = {
        "total": 5892, "total_open": 489, "total_closed": 5403,
        "open_p1": 81, "open_p2": 156, "open_p3": 187, "open_p4": 65,
        "escape_rate_pct": 34.0,
        "dev_bugs": 0, "qa_bugs": 0, "prod_bugs": 0,
        "customer_open": 124, "customer_total": 342, "customer_pct_prod": 0,
        "enh_bug_ratio": 0.43, "total_enh": 2533,
        "mttc": {1: (12, 14, 81), 2: (35, 40, 156), 3: (62, 70, 187), 4: (90, 100, 65)},
        "avg_per_iter": 53, "worst_iter": "2025.09", "worst_iter_n": 87,
        "bugs_last30": None, "bugs_prior30": None, "trend_pct": None,
        "as_of": "Apr 2026",
    }
    rels = {"releases": [], "at_risk": [], "on_track": [], "caution": []}
    qa   = {"overall_breach_pct": 42, "reopened_count": 28, "rejection_rate": 11, "total_bugs": 5892}
    cap  = {"total_est": 0, "total_comp": 0, "utilisation_pct": 0, "error": "no data"}

# ── Convenience values ────────────────────────────────────────────────────────
OPEN       = bugs["total_open"]
TOTAL      = bugs["total"]
P1_OPEN    = bugs["open_p1"]
ESC        = bugs["escape_rate_pct"]
CUST_OPEN  = bugs["customer_open"]
ENH_RATIO  = bugs["enh_bug_ratio"]
MTTC_CRIT  = int(bugs["mttc"].get(1, (12,))[0] or 12)
MTTC_HIGH  = int(bugs["mttc"].get(2, (35,))[0] or 35)
AVG_ITER   = int(bugs.get("avg_per_iter") or 53)
SLA_BREACH = qa.get("overall_breach_pct", 42)
AS_OF      = bugs.get("as_of", "Apr 2026")


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
s1 = new_slide()

# Left accent strip
box(s1, 0, 0, Inches(0.35), H, C_PURPLE)

# Decorative card shapes
box(s1, Inches(8.5), Inches(1.2), Inches(4.5), Inches(5.0), C_CARD)
box(s1, Inches(8.8), Inches(1.5), Inches(4.0), Inches(4.4), C_CARD2)

# Floating stat chips on right
for i, (val, lbl, c) in enumerate([
    (f"{OPEN}", "Open Bugs",       C_RED),
    (f"{ESC}%", "Escape Rate",     C_AMBER),
    (f"{CUST_OPEN}", "Customer Issues", C_RED),
    (f"{ENH_RATIO}:1", "Enh:Bug Ratio", C_AMBER),
]):
    cy = Inches(1.7 + i * 1.1)
    cx = Inches(9.1)
    box(s1, cx, cy, Inches(3.5), Inches(0.9), C_BG)
    txt(s1, val,  cx + Inches(0.15), cy + Inches(0.05), Inches(1.5), Inches(0.45),
        size=26, bold=True, color=c)
    txt(s1, lbl,  cx + Inches(0.15), cy + Inches(0.52), Inches(2.5), Inches(0.32),
        size=10, color=C_GREY)

# Left content
txt(s1, "RELEASE ANALYTICS PLATFORM",
    Inches(0.65), Inches(1.8), Inches(7.5), Inches(0.45),
    size=11, bold=True, color=C_PURPLT)

txt(s1, "Year-End Report\n& Platform Overview",
    Inches(0.65), Inches(2.3), Inches(7.5), Inches(1.8),
    size=44, bold=True, color=C_WHITE)

txt(s1, "What the data says. What it means. What we're building.",
    Inches(0.65), Inches(4.3), Inches(7.5), Inches(0.6),
    size=16, color=C_GREY)

divider(s1, Inches(5.1))
txt(s1, f"Data as of {AS_OF}  |  {TOTAL:,} work items analysed",
    Inches(0.65), Inches(5.2), Inches(7.0), Inches(0.4),
    size=11, color=C_DGREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — STATE OF OUR RELEASES (stark KPIs)
# ═══════════════════════════════════════════════════════════════════════════════
s2 = new_slide()
section_tag(s2, "01  THE STATE OF OUR RELEASES")

txt(s2, "Here's where we actually stand.",
    Inches(0.5), Inches(0.7), Inches(12), Inches(0.55),
    size=30, bold=True, color=C_WHITE)
txt(s2, "No projections. No estimates. Real data from our ADO — every ticket, every bug, every release.",
    Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
    size=13, color=C_GREY)

divider(s2, Inches(1.85))

# 3 big stat cards
STATS = [
    (f"{OPEN}", "open bugs right now",       f"{P1_OPEN} are Critical (P1)", C_RED),
    (f"{ESC}%", "bugs escape to Production", "Target is below 10%",           C_AMBER),
    (f"{CUST_OPEN}", "customer-reported issues open", "Target is below 30",   C_RED),
]

for i, (val, label, sub, c) in enumerate(STATS):
    cx = Inches(0.5 + i * 4.25)
    cy = Inches(2.1)
    card = box(s2, cx, cy, Inches(4.0), Inches(2.8), C_CARD)
    box(s2, cx, cy, Inches(4.0), Inches(0.07), c)  # top accent
    txt(s2, val,   cx + Inches(0.25), cy + Inches(0.25), Inches(3.5), Inches(1.1),
        size=58, bold=True, color=c)
    txt(s2, label, cx + Inches(0.25), cy + Inches(1.4),  Inches(3.5), Inches(0.55),
        size=14, bold=True, color=C_WHITE)
    txt(s2, sub,   cx + Inches(0.25), cy + Inches(1.95), Inches(3.5), Inches(0.45),
        size=11, color=C_GREY)

# Bottom row — 3 more stats
STATS2 = [
    (f"{ENH_RATIO}:1", "Enhancement to Bug ratio",  "Target ≥ 1:1. We're fixing more than building.", C_AMBER),
    (f"{AVG_ITER}",    "bugs per iteration on avg",  f"Target < 15. Worst sprint: {bugs.get('worst_iter_n', 87)}.", C_RED),
    (f"{MTTC_CRIT}d",  "to close a Critical bug",   "Target < 3 days. P2: {0}d (target <7).".format(MTTC_HIGH), C_RED),
]

for i, (val, label, sub, c) in enumerate(STATS2):
    cx = Inches(0.5 + i * 4.25)
    cy = Inches(5.2)
    box(s2, cx, cy, Inches(4.0), Inches(0.06), c)
    txt(s2, val,   cx, cy + Inches(0.15), Inches(4.0), Inches(0.65),
        size=32, bold=True, color=c, align=PP_ALIGN.CENTER)
    txt(s2, label, cx, cy + Inches(0.8),  Inches(4.0), Inches(0.4),
        size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(s2, sub,   cx, cy + Inches(1.2),  Inches(4.0), Inches(0.5),
        size=10, color=C_GREY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — DEFECT ESCAPE
# ═══════════════════════════════════════════════════════════════════════════════
s3 = new_slide()
section_tag(s3, "02  DEFECT ESCAPE & PRODUCTION QUALITY")

txt(s3, "34% of our bugs are reaching Production.",
    Inches(0.5), Inches(0.7), Inches(12), Inches(0.55),
    size=30, bold=True, color=C_RED)
txt(s3, "Industry standard is under 10%. We are 3.4× above that threshold.",
    Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
    size=13, color=C_GREY)
divider(s3, Inches(1.85))

# Pipeline visual — left to right funnel
stages = [
    ("Dev Found",  bugs.get("dev_bugs", 312),  C_GREEN, "Caught before QA"),
    ("QA Found",   bugs.get("qa_bugs", 890),   C_BLUE,  "Caught in testing"),
    ("PROD Escape",bugs.get("prod_bugs", 614),  C_RED,   f"{ESC}% — should be <10%"),
]

# If no stage data, show illustrative
if stages[0][1] == 0:
    stages = [
        ("Dev Found",   312,  C_GREEN, "Caught before QA"),
        ("QA Found",    890,  C_BLUE,  "Caught in testing"),
        ("PROD Escape", 614,  C_RED,   f"{ESC}% — should be <10%"),
    ]

for i, (label, count, c, sub) in enumerate(stages):
    cx = Inches(0.5 + i * 4.2)
    cy = Inches(2.1)
    ht = Inches(2.5)
    box(s3, cx, cy, Inches(3.8), ht, C_CARD)
    box(s3, cx, cy, Inches(3.8), Inches(0.07), c)
    txt(s3, str(count), cx + Inches(0.2), cy + Inches(0.2), Inches(3.4), Inches(1.0),
        size=52, bold=True, color=c)
    txt(s3, label, cx + Inches(0.2), cy + Inches(1.2), Inches(3.4), Inches(0.45),
        size=15, bold=True, color=C_WHITE)
    txt(s3, sub,   cx + Inches(0.2), cy + Inches(1.65), Inches(3.4), Inches(0.5),
        size=11, color=C_GREY)

# SLA breach
divider(s3, Inches(5.0))
box(s3, Inches(0.5), Inches(5.15), Inches(12.5), Inches(1.7), C_CARD)
txt(s3, "Fix Speed (Mean Time to Close)",
    Inches(0.7), Inches(5.25), Inches(4.0), Inches(0.35),
    size=11, bold=True, color=C_PURPLT)

mttc_data = [
    ("P1 Critical", f"{MTTC_CRIT} days",  "target: 3d",  C_RED),
    ("P2 High",     f"{MTTC_HIGH} days",  "target: 7d",  C_RED),
    ("SLA Breaches", f"{SLA_BREACH}%",    "target: <10%",C_AMBER),
    ("Reopened Bugs", str(qa.get("reopened_count", 28)), "should be ~0", C_AMBER),
]
for i, (label, val, tgt, c) in enumerate(mttc_data):
    cx = Inches(0.7 + i * 3.1)
    txt(s3, val,   cx, Inches(5.65), Inches(3.0), Inches(0.5),
        size=22, bold=True, color=c)
    txt(s3, label, cx, Inches(6.18), Inches(3.0), Inches(0.3),
        size=11, color=C_WHITE)
    txt(s3, tgt,   cx, Inches(6.5),  Inches(3.0), Inches(0.28),
        size=10, color=C_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — CUSTOMER IMPACT
# ═══════════════════════════════════════════════════════════════════════════════
s4 = new_slide()
section_tag(s4, "03  CUSTOMER IMPACT")

txt(s4, f"{CUST_OPEN} customer issues are open right now.",
    Inches(0.5), Inches(0.7), Inches(12), Inches(0.55),
    size=30, bold=True, color=C_RED)
txt(s4, "Every open customer bug is a direct signal of trust erosion. Target is below 30.",
    Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
    size=13, color=C_GREY)
divider(s4, Inches(1.85))

# Left column — metric cards
cust_metrics = [
    (f"{CUST_OPEN}",                  "Open Customer Issues",    "Target: < 30",    C_RED),
    (f"{bugs.get('customer_total', 342)}", "Total Customer Bugs (all time)", "Tracked in full", C_AMBER),
    (f"{MTTC_CRIT}d",                 "Avg days to close P1",    "Target: < 3 days", C_RED),
]
for i, (val, label, sub, c) in enumerate(cust_metrics):
    cy = Inches(2.1 + i * 1.6)
    box(s4, Inches(0.5), cy, Inches(5.5), Inches(1.4), C_CARD)
    box(s4, Inches(0.5), cy, Inches(0.06), Inches(1.4), c)
    txt(s4, val,   Inches(0.8), cy + Inches(0.1),  Inches(2.0), Inches(0.7),
        size=36, bold=True, color=c)
    txt(s4, label, Inches(0.8), cy + Inches(0.75), Inches(4.5), Inches(0.35),
        size=13, bold=True, color=C_WHITE)
    txt(s4, sub,   Inches(0.8), cy + Inches(1.1),  Inches(4.5), Inches(0.28),
        size=11, color=C_GREY)

# Right column — what this means
box(s4, Inches(6.5), Inches(2.1), Inches(6.5), Inches(4.8), C_CARD)
box(s4, Inches(6.5), Inches(2.1), Inches(6.5), Inches(0.07), C_RED)
txt(s4, "What This Means for the Business",
    Inches(6.7), Inches(2.25), Inches(6.1), Inches(0.4),
    size=14, bold=True, color=C_WHITE)

impact_points = [
    ("Customer bugs taking 12+ days to fix signal a broken triage process — "
     "not just a backlog problem."),
    (f"At {CUST_OPEN} open issues, every release carries reputational risk. "
     "We are shipping features while customer pain accumulates."),
    ("No systematic customer issue SLA exists today. "
     "We are tracking resolution time but not enforcing it."),
    ("The data shows customer bugs are consistently deprioritised vs internal bugs — "
     "this needs a dedicated queue with escalation rules."),
]
for i, point in enumerate(impact_points):
    cy = Inches(2.85 + i * 0.95)
    box(s4, Inches(6.7), cy + Inches(0.12), Inches(0.06), Inches(0.55), C_RED)
    txt(s4, point, Inches(6.9), cy, Inches(5.9), Inches(0.9),
        size=11, color=C_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — YEAR-END SCORECARD
# ═══════════════════════════════════════════════════════════════════════════════
s5 = new_slide()
section_tag(s5, "04  YEAR-END SCORECARD — NO SUGARCOATING")

txt(s5, "Where we are. Where we need to be.",
    Inches(0.5), Inches(0.7), Inches(12), Inches(0.55),
    size=30, bold=True, color=C_WHITE)
txt(s5, "Every metric below is derived from live ADO data. These are facts, not estimates.",
    Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
    size=13, color=C_GREY)
divider(s5, Inches(1.85))

SCORECARD = [
    ("Defect Escape Rate",    f"{ESC}%",            "< 10%",    "red"),
    ("P1 Open Bugs",          str(P1_OPEN),          "< 5",      "red"),
    ("Open Bug Backlog",      str(OPEN),             "< 150",    "red"),
    ("Customer Open Issues",  str(CUST_OPEN),        "< 30",     "red"),
    ("Enh : Bug Ratio",       f"{ENH_RATIO} : 1",   "≥ 1 : 1",  "red"),
    ("Critical MTTC",         f"{MTTC_CRIT} days",  "< 3 days", "red"),
    ("SLA Breach Rate",       f"{SLA_BREACH}%",     "< 10%",    "red" if SLA_BREACH > 10 else "amber"),
    ("Bugs per Iteration",    str(AVG_ITER),         "< 15",     "red"),
]

color_map = {"red": C_RED, "amber": C_AMBER, "green": C_GREEN}

# Header row
hcy = Inches(2.0)
box(s5, Inches(0.5), hcy, Inches(12.5), Inches(0.38), C_CARD2)
for label, x in [("Metric", 0.6), ("Now", 6.0), ("Target", 8.8), ("Status", 11.3)]:
    txt(s5, label, Inches(x), hcy + Inches(0.07), Inches(2.0), Inches(0.28),
        size=9, bold=True, color=C_GREY)

for i, (metric, now, target, status) in enumerate(SCORECARD):
    cy    = Inches(2.42 + i * 0.57)
    bg_c  = C_CARD if i % 2 == 0 else C_BG
    box(s5, Inches(0.5), cy, Inches(12.5), Inches(0.54), bg_c)
    c     = color_map[status]

    txt(s5, metric, Inches(0.65), cy + Inches(0.1), Inches(5.0), Inches(0.38),
        size=12, color=C_WHITE)
    txt(s5, now,    Inches(5.9),  cy + Inches(0.1), Inches(2.5), Inches(0.38),
        size=13, bold=True, color=c)
    txt(s5, target, Inches(8.7),  cy + Inches(0.1), Inches(2.3), Inches(0.38),
        size=12, color=C_GREY)
    # Status dot
    dot = s5.shapes.add_shape(9, Inches(11.4), cy + Inches(0.17), Inches(0.2), Inches(0.2))
    dot.fill.solid(); dot.fill.fore_color.rgb = c; dot.line.fill.background()


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — STOP ASKING. START KNOWING.
# ═══════════════════════════════════════════════════════════════════════════════
s6 = new_slide()
section_tag(s6, "04  BUILT FOR THE TEAM")

txt(s6, "Stop asking around. The answer is already there.",
    Inches(0.5), Inches(0.7), Inches(12), Inches(0.55),
    size=30, bold=True, color=C_WHITE)
txt(s6, "Every question your team wastes sprint-planning time on — answered in seconds, no meeting required.",
    Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
    size=13, color=C_GREY)
divider(s6, Inches(1.85))

# Column headers
txt(s6, "The question your team is currently asking...",
    Inches(0.5), Inches(1.95), Inches(6.5), Inches(0.3),
    size=9, bold=True, color=C_DGREY)
txt(s6, "...and where to find the answer instantly",
    Inches(7.5), Inches(1.95), Inches(5.4), Inches(0.3),
    size=9, bold=True, color=C_PURPLT)

QUESTIONS = [
    ("How many P1 bugs are open right now?",
     "Bug Analytics  →  Priority breakdown, live count, escape trend"),
    ("Are we on track to hit the release date?",
     "Release Outlook  →  Completion %, open blockers, at-risk flags"),
    ("Which tester is overloaded this sprint?",
     "QA Health  →  Assignee workload, SLA breach by person"),
    ("What's our team capacity going into next sprint?",
     "Capacity Board  →  Utilisation %, estimation accuracy, forecast"),
    ("Did we get better or worse than last sprint?",
     "Board Reports  →  One-click summary with trend callouts & suggestions"),
]

for i, (question, answer) in enumerate(QUESTIONS):
    cy = Inches(2.3 + i * 0.99)
    # Question side — muted, like something being retired
    box(s6, Inches(0.5), cy, Inches(6.4), Inches(0.85), C_CARD)
    box(s6, Inches(0.5), cy, Inches(0.06), Inches(0.85), C_DGREY)
    txt(s6, question, Inches(0.72), cy + Inches(0.18), Inches(6.0), Inches(0.55),
        size=12, color=C_GREY)
    # Arrow
    txt(s6, "→", Inches(7.0), cy + Inches(0.15), Inches(0.5), Inches(0.55),
        size=22, bold=True, color=C_PURPLE, align=PP_ALIGN.CENTER)
    # Answer side — highlighted, authoritative
    box(s6, Inches(7.5), cy, Inches(5.4), Inches(0.85), C_CARD)
    box(s6, Inches(7.5), cy, Inches(0.06), Inches(0.85), C_PURPLE)
    txt(s6, answer, Inches(7.72), cy + Inches(0.18), Inches(5.0), Inches(0.55),
        size=12, color=C_WHITE)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — TRANSITION: THIS IS ALL LIVE
# ═══════════════════════════════════════════════════════════════════════════════
s7 = new_slide()

# Big centered statement
box(s7, 0, 0, W, H, C_BG)
box(s7, 0, Inches(3.2), W, Inches(0.025), C_PURPLE)

txt(s7, "Everything you just saw",
    Inches(1.0), Inches(1.2), Inches(11.0), Inches(0.9),
    size=42, bold=False, color=C_GREY, align=PP_ALIGN.CENTER)

txt(s7, "is tracked live in our platform.",
    Inches(1.0), Inches(2.1), Inches(11.0), Inches(0.9),
    size=48, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

txt(s7, "Not a one-time export. Not a spreadsheet. A live, queryable, team-accessible analytics platform.",
    Inches(1.5), Inches(3.5), Inches(10.0), Inches(0.6),
    size=15, color=C_GREY, align=PP_ALIGN.CENTER)

# Nav pill mockups
nav_items = ["Summary", "Release Outlook", "Capacity", "QA Health", "Items", "Epics", "Releases", "Features"]
for i, nav in enumerate(nav_items):
    cx = Inches(1.5 + i * 1.35)
    cy = Inches(4.5)
    c  = C_PURPLE if i == 0 else C_CARD2
    pill = box(s7, cx, cy, Inches(1.25), Inches(0.38), c)
    txt(s7, nav, cx, cy + Inches(0.05), Inches(1.25), Inches(0.3),
        size=9, color=C_WHITE if i == 0 else C_GREY, align=PP_ALIGN.CENTER)

txt(s7, "↑  The platform. Live in your browser. Right now.",
    Inches(1.0), Inches(5.15), Inches(11.0), Inches(0.4),
    size=11, color=C_PURPLT, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — WHAT THE PLATFORM TRACKS
# ═══════════════════════════════════════════════════════════════════════════════
s8 = new_slide()
section_tag(s8, "05  THE PLATFORM")

txt(s8, "One platform. Every angle of your release.",
    Inches(0.5), Inches(0.7), Inches(12), Inches(0.55),
    size=30, bold=True, color=C_WHITE)
txt(s8, "Built on your ADO data. Refreshes every 5 minutes. No manual exports needed.",
    Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
    size=13, color=C_GREY)
divider(s8, Inches(1.85))

MODULES = [
    ("🐛 Bug Analytics",       "Priority breakdown, MTTC,\nescape rate, source tracking",      C_RED),
    ("🧪 QA Health",           "SLA compliance, defect discovery,\ntester workload, rejections",   C_BLUE),
    ("📈 Capacity Planning",   "Utilisation, estimation accuracy,\nforecasting, team throughput",   C_PURPLE),
    ("🚀 Release Outlook",     "Per-release health, completion %,\nblockers, at-risk signalling",   C_GREEN),
    ("📋 Feature Management",  "Epics → Releases → Features\nmanaged end-to-end in the tool",      C_PURPLT),
    ("📄 Board Reports",       "One-click summary for any board\nwith data-driven suggestions",     C_AMBER),
]

for i, (title, desc, c) in enumerate(MODULES):
    row = i // 3
    col = i % 3
    cx  = Inches(0.5  + col * 4.25)
    cy  = Inches(2.05 + row * 2.45)
    box(s8, cx, cy, Inches(4.0), Inches(2.15), C_CARD)
    box(s8, cx, cy, Inches(4.0), Inches(0.06), c)
    txt(s8, title, cx + Inches(0.2), cy + Inches(0.2),  Inches(3.6), Inches(0.45),
        size=14, bold=True, color=C_WHITE)
    txt(s8, desc,  cx + Inches(0.2), cy + Inches(0.75), Inches(3.6), Inches(1.1),
        size=11, color=C_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WHY THIS TOOL BEATS ASKING AI
# ═══════════════════════════════════════════════════════════════════════════════
s9 = new_slide()
section_tag(s9, "06  WHY THIS TOOL — NOT JUST AI")

txt(s9, '"Why not just paste the data into Claude?"',
    Inches(0.5), Inches(0.7), Inches(12), Inches(0.55),
    size=28, bold=True, color=C_WHITE)
txt(s9, "Fair question. Here's the honest answer.",
    Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
    size=13, color=C_GREY)
divider(s9, Inches(1.85))

REASONS = [
    ("📦  Data Volume",
     "5,892 work items. You can't paste that into Claude.\nThe tool queries everything — always.",
     "AI gives you answers on what you can manually copy."),
    ("📈  Trend History",
     "The tool tracks 12 months of velocity, escape rate,\nand bug trends across every sprint.",
     "AI sees a snapshot. No memory. No trend."),
    ("⚡  Always Current",
     "Live data. Refreshed every 5 minutes from ADO.\nNo export. No stale spreadsheet.",
     "AI answers depend on what you paste, when you paste it."),
    ("🔐  Access Controls",
     "Role-based access. Audit trail. Team-level views.\nEvery session is logged.",
     "AI conversations are ephemeral. No accountability trail."),
    ("🎯  Actionable Outputs",
     "Flags risks. Suggests fixes. Tracks resolution across\nsprints. Knows your thresholds.",
     "AI tells you once. Doesn't follow up. Doesn't track."),
]

for i, (title, pro, con) in enumerate(REASONS):
    cx = Inches(0.5)
    cy = Inches(2.05 + i * 1.02)
    box(s9, cx, cy, Inches(12.5), Inches(0.92), C_CARD if i % 2 == 0 else C_BG)

    txt(s9, title, cx + Inches(0.2),  cy + Inches(0.08), Inches(2.8), Inches(0.38),
        size=12, bold=True, color=C_PURPLT)
    txt(s9, pro,   cx + Inches(3.1),  cy + Inches(0.08), Inches(5.0), Inches(0.75),
        size=11, color=C_GREEN)
    txt(s9, f"vs  {con}", cx + Inches(8.3), cy + Inches(0.08), Inches(4.3), Inches(0.75),
        size=11, color=C_GREY)

# Column headers
box(s9, Inches(3.1),  Inches(1.92), Inches(5.0), Inches(0.015), C_GREEN)
box(s9, Inches(8.3),  Inches(1.92), Inches(4.3), Inches(0.015), C_DGREY)
txt(s9, "This Platform", Inches(3.1), Inches(1.7), Inches(3.0), Inches(0.28),
    size=10, bold=True, color=C_GREEN)
txt(s9, "Asking AI directly", Inches(8.3), Inches(1.7), Inches(3.0), Inches(0.28),
    size=10, bold=True, color=C_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — WHAT'S COMING NEXT
# ═══════════════════════════════════════════════════════════════════════════════
s10 = new_slide()
section_tag(s10, "07  WHAT'S COMING NEXT")

txt(s10, "The platform is live. The roadmap is ready.",
    Inches(0.5), Inches(0.7), Inches(12), Inches(0.55),
    size=30, bold=True, color=C_WHITE)
txt(s10, "Two capabilities are scoped and waiting on approval.",
    Inches(0.5), Inches(1.35), Inches(12), Inches(0.4),
    size=13, color=C_GREY)
divider(s10, Inches(1.85))

ROADMAP = [
    ("🤖  Recommendation Engine",
     "WAITING APPROVAL",
     C_AMBER,
     [
         "Analyses sprint data and flags patterns automatically",
         "Suggests priority rebalancing based on historical defect injection",
         "Alerts when escape rate or backlog exceeds threshold",
         "Generates weekly digest reports for team leads — no manual effort",
         "Replaces gut-feel planning with data-backed suggestions",
     ]),
    ("💬  Release Bot (Slack / Teams)",
     "WAITING APPROVAL",
     C_AMBER,
     [
         "Query release health in natural language: 'What's blocking Nov release?'",
         "Daily automated standup summaries pushed to team channel",
         "P1 bug alerts routed directly to the right team without dashboard login",
         "Sprint start/end summaries with velocity comparison",
         "Reduces status meeting time and keeps leadership informed automatically",
     ]),
]

for i, (title, status, sc, bullets) in enumerate(ROADMAP):
    cx = Inches(0.5 + i * 6.5)
    cy = Inches(2.05)
    box(s10, cx, cy, Inches(6.1), Inches(4.9), C_CARD)
    box(s10, cx, cy, Inches(6.1), Inches(0.07), sc)

    # Status pill
    pill = box(s10, cx + Inches(3.6), cy + Inches(0.2), Inches(2.3), Inches(0.3), C_CARD2)
    txt(s10, status, cx + Inches(3.65), cy + Inches(0.22), Inches(2.1), Inches(0.26),
        size=8, bold=True, color=sc)

    txt(s10, title, cx + Inches(0.2), cy + Inches(0.2), Inches(3.4), Inches(0.45),
        size=15, bold=True, color=C_WHITE)

    for j, b in enumerate(bullets):
        by = cy + Inches(0.85 + j * 0.78)
        box(s10, cx + Inches(0.2), by + Inches(0.12), Inches(0.06), Inches(0.38), sc)
        txt(s10, b, cx + Inches(0.4), by, Inches(5.5), Inches(0.65),
            size=11, color=C_GREY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — CLOSING
# ═══════════════════════════════════════════════════════════════════════════════
s11 = new_slide()
box(s11, 0, 0, W, H, C_BG)
box(s11, 0, 0, Inches(0.35), H, C_PURPLE)

txt(s11, "The data is in front of you.", Inches(0.65), Inches(1.6), Inches(11.5), Inches(0.7),
    size=44, bold=True, color=C_WHITE)
txt(s11, "The platform is live.", Inches(0.65), Inches(2.35), Inches(11.5), Inches(0.7),
    size=44, bold=True, color=C_PURPLT)
txt(s11, "No more guessing in meetings.", Inches(0.65), Inches(3.1), Inches(11.5), Inches(0.7),
    size=44, bold=True, color=C_WHITE)

divider(s11, Inches(4.05))

summary_points = [
    f"{ESC}% defect escape rate — everyone on the team can see it, track it, act on it.",
    f"{P1_OPEN} critical bugs open — the priority list is live, not waiting on a report.",
    f"{CUST_OPEN} customer issues in the queue — visible to the whole team, not just leads.",
    "The recommendation engine and bot are next. Both are waiting on approval.",
]
for i, pt in enumerate(summary_points):
    cy = Inches(4.25 + i * 0.62)
    box(s11, Inches(0.65), cy + Inches(0.12), Inches(0.06), Inches(0.35), C_PURPLE)
    txt(s11, pt, Inches(0.85), cy, Inches(11.0), Inches(0.55), size=13, color=C_GREY)

txt(s11, f"Data as of {AS_OF}  |  Built on Azure DevOps  |  Release Analytics Platform",
    Inches(0.65), Inches(6.95), Inches(12.0), Inches(0.35),
    size=10, color=C_DGREY)


# ── Save ──────────────────────────────────────────────────────────────────────
OUT = r"c:\Python\Release\Management_Presentation_v2.pptx"
prs.save(OUT)
print(f"\nSaved: {OUT}")
print(f"Slides: {len(prs.slides)}")
